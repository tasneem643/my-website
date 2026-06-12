from flask import Flask, render_template, request, redirect, url_for, session, flash, send_from_directory, abort
app = Flask(__name__)  # تم نقل السطر هنا في البداية لتفادي خطأ التعريف

import pymysql
import db  # طبقة الوصول المزدوجة (SQLite / MySQL)
import os
import time
import re
import secrets
import string
import logging
from datetime import timedelta
from collections import defaultdict
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash

# مكاتب المعالجة وقراءة الملفات
import pypdf
import pytesseract
from docx import Document
from pptx import Presentation   # مكتبة لقراءة PowerPoint
from PIL import Image

# حل ذكي لتفادي قفل الـ DLL محلياً بسبب إصدار بايثون الحديث
try:
    import pymupdf as fitz
except ImportError:
    try:
        import fitz
    except Exception:
        fitz = None  # يتخطى الخطأ محلياً حتى يشتغل السيرفر فوراً

# -----------------------------
# فلتر تنظيف النصوص
# -----------------------------
def clean_text(text):
    import re
    # إبقاء العربية + الإنجليزية + الأرقام + بعض الرموز الأساسية
    text = re.sub(r'[^\u0600-\u06FFa-zA-Z0-9\s.,!?؟]', '', text)
    # إزالة التكرار المبالغ فيه للأحرف
    text = re.sub(r'(.)\1{2,}', r'\1', text)
    return text.strip()

# -----------------------------
# دالة استخراج النص من الملفات
# -----------------------------
def extract_text_from_file(file_path):
    text = ""
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            try:
                # محاولة القراءة كنص
                reader = pypdf.PdfReader(file_path)
                for page in reader.pages:
                    text += page.extract_text() or ""
            except Exception:
                # إذا فشلت القراءة، نستخدم OCR لو كانت المكتبة تعمل
                if fitz:
                    doc = fitz.open(file_path)
                    for page in doc:
                        pix = page.get_pixmap()
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        text += pytesseract.image_to_string(img, lang="ara+eng")
                else:
                    text = "⚠️ ميزة OCR معطلة محلياً بسبب البيئة."

        elif ext == ".docx":
            doc = Document(file_path)
            for para in doc.paragraphs:
                if para.text.strip():
                    text += para.text + "\n"

        elif ext == ".pptx":
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"

        else:
            text = "⚠️ نوع الملف غير مدعوم حالياً."

    except Exception as e:
        text = f"⚠️ الملف غير صالح للقراءة: {e}"

    return clean_text(text)

# -----------------------------
# دالة التلخيص
# -----------------------------
def summarize_handout(file_name):
    full_path = os.path.join(app.config['UPLOAD_FOLDER'], file_name)
    content = extract_text_from_file(full_path)

    if not content or len(content.strip()) < 20:
        return "لا يوجد محتوى كافي للتلخيص."

    import re
    sentences = re.split(r'[.!?؟]', content)
    summary = " ".join(sentences[:3])

    return summary.strip()

# إعدادات المجلد وقاعدة البيانات
DATABASE = 'musaid_ist.db'
BASE_DIR = os.path.abspath(os.path.dirname(__file__))
UPLOAD_FOLDER = os.path.join(BASE_DIR, 'uploads')
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

# =====================================================================
# الأمان: مفتاح الجلسة، كعكات آمنة، CSRF، تحديد المحاولات، سجل التدقيق
# =====================================================================

try:
    from dotenv import load_dotenv
    load_dotenv(os.path.join(BASE_DIR, '.env'))
except Exception:
    pass

def _load_secret_key():
    env_key = os.environ.get('MUSAID_SECRET_KEY')
    if env_key:
        return env_key
    key_file = os.path.join(BASE_DIR, '.secret_key')
    try:
        if os.path.exists(key_file):
            with open(key_file, 'r', encoding='utf-8') as fh:
                saved = fh.read().strip()
                if saved:
                    return saved
        generated = secrets.token_hex(32)
        with open(key_file, 'w', encoding='utf-8') as fh:
            fh.write(generated)
        return generated
    except Exception:
        return secrets.token_hex(32)

app.secret_key = _load_secret_key()

ENV = os.environ.get('MUSAID_ENV', 'production').strip().lower()
IS_PRODUCTION = ENV not in ('development', 'dev', 'local')

def _env_bool(name, default=False):
    val = os.environ.get(name)
    if val is None:
        return default
    return val.strip().lower() in ('1', 'true', 'yes', 'on')

DEBUG = False if IS_PRODUCTION else _env_bool('MUSAID_DEBUG', True)

app.config.update(
    DEBUG=DEBUG,
    ENV=ENV,
    SESSION_COOKIE_HTTPONLY=True,
    SESSION_COOKIE_SAMESITE='Lax',
    SESSION_COOKIE_SECURE=_env_bool('MUSAID_COOKIE_SECURE', IS_PRODUCTION),
    PERMANENT_SESSION_LIFETIME=timedelta(hours=8),
    MAX_CONTENT_LENGTH=50 * 1024 * 1024,
)

# --- سجل تدقيق المصادقة ---
LOG_DIR = os.path.join(BASE_DIR, 'logs')
try:
    os.makedirs(LOG_DIR, exist_ok=True)
except Exception:
    pass
audit_logger = logging.getLogger('musaid.auth')
audit_logger.setLevel(logging.INFO)
if not audit_logger.handlers:
    try:
        _audit_handler = logging.FileHandler(
            os.path.join(LOG_DIR, 'auth.log'), encoding='utf-8')
        _audit_handler.setFormatter(
            logging.Formatter('%(asctime)s %(message)s'))
        audit_logger.addHandler(_audit_handler)
    except Exception:
        pass

def _client_ip():
    try:
        fwd = request.headers.get('X-Forwarded-For')
        if fwd:
            return fwd.split(',')[0].strip()
        return request.remote_addr or 'unknown'
    except Exception:
        return 'unknown'

def audit_log(event, detail=''):
    try:
        user = session.get('user_id', '-')
    except Exception:
        user = '-'
    try:
        audit_logger.info(f'event={event} ip={_client_ip()} user={user} {detail}'.strip())
    except Exception:
        pass

# --- حماية CSRF ---
def get_csrf_token():
    token = session.get('_csrf_token')
    if not token:
        token = secrets.token_hex(32)
        session['_csrf_token'] = token
    return token

@app.context_processor
def _inject_csrf():
    return {'csrf_token': get_csrf_token}

@app.before_request
def _csrf_protect():
    if request.method in ('POST', 'PUT', 'PATCH', 'DELETE'):
        sent = (request.form.get('csrf_token')
                or request.headers.get('X-CSRFToken'))
        stored = session.get('_csrf_token')
        if not stored or not sent or not secrets.compare_digest(str(stored), str(sent)):
            audit_log('csrf_failure', f'path={request.path}')
            abort(400)

# --- تحديد محاولات الدخول ---
RATE_LIMIT_WINDOW = 300      
RATE_LIMIT_MAX_FAILS = 5     
LOCKOUT_SECONDS = 900        
_login_failures = defaultdict(list)   
_login_lockouts = {}                  

def _rl_key():
    return _client_ip()

def login_lock_remaining(key):
    until = _login_lockouts.get(key)
    if until and time.time() < until:
        return int(until - time.time())
    if until:
        _login_lockouts.pop(key, None)
    return 0

def record_login_failure(key):
    now = time.time()
    fails = [t for t in _login_failures[key] if now - t < RATE_LIMIT_WINDOW]
    fails.append(now)
    _login_failures[key] = fails
    if len(fails) >= RATE_LIMIT_MAX_FAILS:
        _login_lockouts[key] = now + LOCKOUT_SECONDS
        _login_failures[key] = []
        return True
    return False

def clear_login_failures(key):
    _login_failures.pop(key, None)
    _login_lockouts.pop(key, None)

_EMAIL_RE = re.compile(r'^[^@\s]+@[^@\s]+\.[^@\s]+$')

def validate_email(email):
    return bool(email) and len(email) <= 120 and bool(_EMAIL_RE.match(email))

def clean_name(name):
    return (name or '').strip()[:120]

# --- دالة الاتصال الذكية بالبيانات ---
def get_db_connection():
    if os.environ.get('MUSAID_ENV') == 'production' or os.environ.get('DATABASE_URL'):
        return pymysql.connect(
            host=os.environ.get('DB_HOST', 'localhost'),
            user=os.environ.get('DB_USER', 'root'),
            password=os.environ.get('DB_PASSWORD', ''),
            database=os.environ.get('DB_NAME', 'musaid_ist'),
            cursorclass=pymysql.cursors.DictCursor
        )
    else:
        return db.connect(DATABASE)

def ensure_schema():
    conn = get_db_connection()
    try:
        if os.environ.get('MUSAID_ENV') == 'production' or db.is_mysql():
            cursor = conn.cursor()
            cursor.execute(
                "SELECT COLUMN_NAME FROM information_schema.COLUMNS "
                "WHERE TABLE_SCHEMA = DATABASE() AND TABLE_NAME = 'handouts'")
            existing = {row['COLUMN_NAME'] for row in cursor.fetchall()}
            for col in ('view_count', 'download_count'):
                if col not in existing:
                    cursor.execute(f'ALTER TABLE handouts ADD COLUMN {col} INT DEFAULT 0')
        else:
            cols = {row['name'] for row in conn.execute('PRAGMA table_info(handouts)')}
            for col in ('view_count', 'download_count'):
                if col not in cols:
                    conn.execute(f'ALTER TABLE handouts ADD COLUMN {col} INTEGER DEFAULT 0')
        conn.commit()
    finally:
        conn.close()

try:
    ensure_schema()
except Exception as _e:
    print(f'⚠️  ensure_schema تخطّى التهيئة: {_e}')

# -----------------------------
# أمان كلمات المرور
# -----------------------------
_HASH_PREFIXES = ('pbkdf2:', 'scrypt:', 'argon2')

def is_hashed(value):
    return bool(value) and str(value).startswith(_HASH_PREFIXES)

def verify_password(stored, candidate):
    if not stored:
        return False
    if is_hashed(stored):
        try:
            return check_password_hash(stored, candidate)
        except Exception:
            return False
    return secrets.compare_digest(str(stored), str(candidate))

def hash_password(plain):
    return generate_password_hash(plain, method='pbkdf2:sha256', salt_length=16)

def validate_password(pw):
    if not pw or len(pw) < 8:
        return False, 'يجب أن تتكون كلمة المرور من 8 خانات على الأقل.'
    if not re.search(r'[A-Za-z]', pw):
        return False, 'يجب أن تحتوي كلمة المرور على حرف واحد على الأقل.'
    if not re.search(r'\d', pw):
        return False, 'يجب أن تحتوي كلمة المرور على رقم واحد على الأقل.'
    return True, ''

def generate_strong_password(length=12):
    alphabet = string.ascii_letters + string.digits + '!@#$%&*?'
    while True:
        pw = ''.join(secrets.choice(alphabet) for _ in range(length))
        if (re.search(r'[a-z]', pw) and re.search(r'[A-Z]', pw)
                and re.search(r'\d', pw) and re.search(r'[!@#$%&*?]', pw)):
            return pw

# --- دالة التحميل والمعاينة ---
@app.route('/download/<filename>')
def uploaded_file(filename):
    try:
        column = 'download_count' if request.args.get('dl') else 'view_count'
        conn = get_db_connection()
        cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
        
        query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
        cursor.execute(
            f'UPDATE handouts SET {column} = COALESCE({column}, 0) + 1 WHERE file_path = {query_placeholder}',
            (filename,))
        conn.commit()
        conn.close()
    except Exception:
        pass
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename, as_attachment=False)

# --- خدمة ملفات PWA ---
@app.route('/service-worker.js')
def service_worker():
    return send_from_directory(BASE_DIR, 'service-worker.js', mimetype='application/javascript')

@app.route('/manifest_student.json')
def manifest_student():
    return send_from_directory(BASE_DIR, 'manifest_student.json', mimetype='application/manifest+json')

@app.route('/manifest_teacher.json')
def manifest_teacher():
    return send_from_directory(BASE_DIR, 'manifest_teacher.json', mimetype='application/manifest+json')

# --- 1. واجهة الطالب ---
@app.route('/')
def index():
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    
    cursor.execute('SELECT * FROM departments')
    depts = cursor.fetchall()

    cursor.execute('SELECT COUNT(*) FROM handouts')
    h_count = cursor.fetchone()
    cursor.execute('SELECT COUNT(*) FROM subjects')
    s_count = cursor.fetchone()
    cursor.execute("SELECT COUNT(*) FROM teachers WHERE email != 'admin@musaid.edu.ly'")
    t_count = cursor.fetchone()
    cursor.execute('SELECT COUNT(*) FROM departments')
    d_count = cursor.fetchone()

    stats = {
        'handouts': h_count[0] if isinstance(h_count, tuple) else h_count['COUNT(*)'],
        'subjects': s_count[0] if isinstance(s_count, tuple) else s_count['COUNT(*)'],
        'teachers': t_count[0] if isinstance(t_count, tuple) else t_count["COUNT(*)"],
        'departments': d_count[0] if isinstance(d_count, tuple) else d_count['COUNT(*)'],
    }

    cursor.execute('''
        SELECT h.title, h.file_path, h.upload_date, h.dept_id, h.semester,
               COALESCE(h.view_count, 0) AS view_count,
               COALESCE(h.download_count, 0) AS download_count,
               s.subject_name, t.full_name AS teacher_name, d.dept_name
        FROM handouts h
        JOIN subjects s ON h.subject_id = s.id
        JOIN teachers t ON h.teacher_id = t.id
        JOIN departments d ON h.dept_id = d.id
        ORDER BY h.id DESC LIMIT 6
    ''')
    recent = cursor.fetchall()

    conn.close()
    return render_template('index.html', depts=depts, stats=stats, recent=recent)
    
@app.route('/search')
def search():
    dept_id = request.args.get('dept')
    semester = request.args.get('semester')
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    query = f'''
        SELECT h.*, s.subject_name, t.full_name as teacher_name
        FROM handouts h
        JOIN subjects s ON h.subject_id = s.id
        JOIN teachers t ON h.teacher_id = t.id
        WHERE h.dept_id = {query_placeholder} AND h.semester = {query_placeholder}
    '''
    cursor.execute(query, (dept_id, semester))
    results = cursor.fetchall()
    
    cursor.execute(f'SELECT dept_name FROM departments WHERE id = {query_placeholder}', (dept_id,))
    dept_name_row = cursor.fetchone()
    dept_name = dept_name_row['dept_name'] if dept_name_row else "غير معروف"
    
    processed_results = []
    for row in results:
        item = dict(row)
        summary = summarize_handout(item.get('file_path', ''))
        item['flash_summary'] = summary
        processed_results.append(item)
    conn.close()
    
    return render_template('results.html', results=processed_results, dept_name=dept_name, semester=semester)

# --- 2. نظام تسجيل الدخول ---
@app.route('/login', methods=['GET', 'POST'])
def login():
    if request.method == 'POST':
        key = _rl_key()

        remaining = login_lock_remaining(key)
        if remaining > 0:
            audit_log('login_locked', f'email={request.form.get("email","")!r}')
            flash(f'🔒 تم تجاوز عدد المحاولات المسموح. حاول مجدداً بعد {remaining // 60 + 1} دقيقة.')
            return render_template('login.html')

        email = (request.form.get('email') or '').strip()
        password = request.form.get('password') or ''

        conn = get_db_connection()
        cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
        
        query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
        cursor.execute(f'SELECT * FROM teachers WHERE email = {query_placeholder}', (email,))
        user = cursor.fetchone()

        if user and verify_password(user['password'], password):
            if not is_hashed(user['password']):
                cursor.execute(f'UPDATE teachers SET password = {query_placeholder} WHERE id = {query_placeholder}',
                             (hash_password(password), user['id']))
                conn.commit()
            conn.close()

            clear_login_failures(key)
            old_csrf = session.get('_csrf_token')
            session.clear()
            if old_csrf:
                session['_csrf_token'] = old_csrf
            session['user_id'] = user['id']
            session['user_name'] = user['full_name']
            session['role'] = 'admin' if email == 'admin@musaid.edu.ly' else 'teacher'
            audit_log('login_success', f'role={session["role"]} email={email!r}')

            if session['role'] == 'admin':
                return redirect(url_for('admin_dashboard'))
            else:
                return redirect(url_for('teacher_dashboard'))
        else:
            conn.close()
            locked_now = record_login_failure(key)
            audit_log('login_failure', f'email={email!r} locked={locked_now}')
            if locked_now:
                flash(f'🔒 تم تجاوز عدد المحاولات المسموح. الحساب مقفل مؤقتاً لمدة {LOCKOUT_SECONDS // 60} دقيقة.')
            else:
                flash('خطأ في البريد الإلكتروني أو كلمة المرور')

    return render_template('login.html')

# --- 3. لوحة تحكم الأستاذ ---
@app.route('/teacher')
def teacher_dashboard():
    if 'user_id' not in session or session['role'] != 'teacher':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    
    cursor.execute('SELECT id, dept_name FROM departments')
    depts = cursor.fetchall()
    depts_list = [dict(row) for row in depts]
    
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    cursor.execute(f'''
        SELECT h.*, s.subject_name 
        FROM handouts h 
        JOIN subjects s ON h.subject_id = s.id 
        WHERE h.teacher_id = {query_placeholder} 
        ORDER BY h.id DESC
    ''', (session['user_id'],))
    my_handouts = cursor.fetchall()
    my_handouts_list = [dict(row) for row in my_handouts]

    my_count = len(my_handouts_list)

    top_subject = None
    if my_count > 0:
        subject_counts = {}
        for h in my_handouts_list:
            subject_counts[h['subject_name']] = subject_counts.get(h['subject_name'], 0) + 1
        top_subject = max(subject_counts, key=subject_counts.get)

    cursor.execute("SELECT COUNT(*) FROM handouts")
    t_count = cursor.fetchone()
    total_count = t_count[0] if isinstance(t_count, tuple) else t_count["COUNT(*)"]

    participation = round((my_count / total_count) * 100, 1) if total_count > 0 else 0

    ai_message = None
    if my_count == 0:
        ai_message = "🔔 تنبيه ذكي: لم تقم برفع أي مذكرة حتى الآن هذا الفصل."

    conn.close()
    
    return render_template('teacher_dashboard.html', 
                           name=session['user_name'], 
                           depts=depts_list,
                           my_handouts=my_handouts_list,
                           my_count=my_count,
                           top_subject=top_subject,
                           participation=participation,
                           ai_alert=ai_message)

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'user_id' not in session or session['role'] != 'teacher':
        return redirect(url_for('login'))

    dept_id = request.form.get('dept_id')
    semester = request.form.get('semester')
    subject_id = request.form.get('subject_id')
    title = request.form.get('title')
    notes = request.form.get('notes')
    files = request.files.getlist('files[]')

    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    try:
        cursor.execute(f"""
            SELECT COUNT(*) FROM handouts 
            WHERE subject_id = {query_placeholder} AND LOWER(title) = LOWER({query_placeholder})
        """, (subject_id, title))
        dup_row = cursor.fetchone()
        duplicate = dup_row[0] if isinstance(dup_row, tuple) else dup_row["COUNT(*)"]

        if duplicate > 0:
            conn.close()
            flash("⚠️ تنبيه ذكي: هناك مذكرة مشابهة مرفوعة مسبقًا لهذه المادة.")
            return redirect(url_for('teacher_dashboard'))

        for file in files:
            if file and file.filename != '':
                original_filename = file.filename
                ext = os.path.splitext(original_filename)[1]
                base_name = secure_filename(os.path.splitext(original_filename)[0])
                unique_filename = f"{int(time.time())}_{base_name}{ext}"
                file.save(os.path.join(app.config['UPLOAD_FOLDER'], unique_filename))

                cursor.execute(f'''
                    INSERT INTO handouts (teacher_id, subject_id, dept_id, semester, title, notes, file_path)
                    VALUES ({query_placeholder}, {query_placeholder}, {query_placeholder}, {query_placeholder}, {query_placeholder}, {query_placeholder}, {query_placeholder})
                ''', (session['user_id'], subject_id, dept_id, semester, title, notes, unique_filename))
        
        conn.commit()
        flash('✅ تم رفع الملفات بنجاح!')
    except Exception as e:
        print(f"Error: {e}")
        flash('❌ حدث خطأ أثناء الرفع!')
    finally:
        conn.close()
    
    return redirect(url_for('teacher_dashboard'))

# --- 4. لوحة تحكم المدير ---
@app.route('/admin')
def admin_dashboard():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    cursor.execute("SELECT * FROM teachers WHERE email != 'admin@musaid.edu.ly' ORDER BY id DESC")
    teachers = cursor.fetchall()
    conn.close()
    return render_template('admin/dashboard.html', teachers=teachers)

@app.route('/admin/add_teacher', methods=['POST'])
def add_teacher():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    full_name = clean_name(request.form.get('name'))
    email = (request.form.get('email') or '').strip().lower()
    password = request.form.get('password') or ''

    if not full_name:
        flash('⚠️ يرجى إدخال اسم الأستاذ.')
        return redirect(url_for('admin_dashboard'))
    if not validate_email(email):
        flash('⚠️ البريد الإلكتروني غير صالح.')
        return redirect(url_for('admin_dashboard'))

    ok, _ = validate_password(password)
    if not ok:
        password = generate_strong_password()

    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    try:
        cursor.execute(f"INSERT INTO teachers (full_name, email, password) VALUES ({query_placeholder}, {query_placeholder}, {query_placeholder})",
                     (full_name, email, hash_password(password)))
        conn.commit()
        audit_log('teacher_add', f'email={email!r}')
        flash(f'✅ تم تسجيل الأستاذ بنجاح! بيانات الدخول (احفظها الآن): {email} / {password}')
    except Exception as e:
        flash('خطأ: قد يكون البريد مسجلاً مسبقاً!')
    finally:
        conn.close()
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/edit_teacher', methods=['POST'])
def edit_teacher():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    t_id = request.form.get('id')
    name = clean_name(request.form.get('name'))
    email = (request.form.get('email') or '').strip().lower()
    password = (request.form.get('password') or '').strip()

    if not name or not validate_email(email):
        flash('⚠️ الاسم أو البريد الإلكتروني غير صالح.')
        return redirect(url_for('admin_dashboard'))

    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    if password:
        ok, msg = validate_password(password)
        if not ok:
            conn.close()
            flash('⚠️ ' + msg)
            return redirect(url_for('admin_dashboard'))
        cursor.execute(f"UPDATE teachers SET full_name = {query_placeholder}, email = {query_placeholder}, password = {query_placeholder} WHERE id = {query_placeholder}",
                     (name, email, hash_password(password), t_id))
    else:
        cursor.execute(f"UPDATE teachers SET full_name = {query_placeholder}, email = {query_placeholder} WHERE id = {query_placeholder}",
                     (name, email, t_id))
    conn.commit()
    conn.close()
    audit_log('teacher_edit', f'id={t_id} pw_changed={bool(password)}')
    flash('تم تحديث بيانات الأستاذ بنجاح')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/subjects')
def admin_subjects():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    
    cursor.execute("SELECT * FROM departments WHERE id IN (1, 2)")
    depts = cursor.fetchall()
    
    query = '''
        SELECT s.id, s.subject_name, cs.dept_id, cs.semester 
        FROM subjects s
        JOIN course_structure cs ON s.id = cs.subject_id
    '''
    cursor.execute(query)
    subjects = cursor.fetchall()
    conn.close()
    
    ai_message = None
    if len(subjects) == 0:
        ai_message = "🔔 تنبيه ذكي: لا توجد أي مواد مسجلة حتى الآن."
    else:
        if len(subjects) > len(depts) * 10:
            ai_message = "🔔 تنبيه ذكي: هناك عدد كبير من المواد مقارنة بالأقسام، تحقق من التوزيع."
    
    return render_template('admin/subjects.html', depts=depts, subjects=subjects, ai_alert=ai_message)

@app.route('/admin/delete_teacher/<int:id>')
def delete_teacher(id):
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    cursor.execute(f"DELETE FROM teachers WHERE id = {query_placeholder}", (id,))
    conn.commit()
    conn.close()
    audit_log('teacher_delete', f'id={id}')
    flash('تم حذف حساب الأستاذ بنجاح')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/reset_password/<int:id>')
def reset_password(id):
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))

    new_password = generate_strong_password()
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    cursor.execute(f"SELECT email FROM teachers WHERE id = {query_placeholder}", (id,))
    row = cursor.fetchone()
    if row:
        cursor.execute(f"UPDATE teachers SET password = {query_placeholder} WHERE id = {query_placeholder}",
                     (hash_password(new_password), id))
        conn.commit()
        conn.close()
        audit_log('password_reset', f'id={id} email={row["email"]!r}')
        flash(f'🔑 تمت إعادة تعيين كلمة المرور لـ {row["email"]} (احفظها الآن): {new_password}')
    else:
        conn.close()
        flash('⚠️ لم يتم العثور على الحساب.')
    return redirect(url_for('admin_dashboard'))

@app.route('/admin/add_subject', methods=['POST'])
def add_subject():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
        
    name = request.form.get('subject_name')
    dept_id = request.form.get('dept_id')
    semester = request.form.get('semester')
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    try:
        cursor.execute(f"INSERT INTO subjects (subject_name) VALUES ({query_placeholder})", (name,))
        subject_id = cursor.lastrowid
        
        cursor.execute(f'''
            INSERT INTO course_structure (subject_id, dept_id, semester) 
            VALUES ({query_placeholder}, {query_placeholder}, {query_placeholder})
        ''', (subject_id, dept_id, semester))
        
        conn.commit()
        flash('تمت إضافة المادة وربطها بالقسم بنجاح!')
    except Exception as e:
        flash('حدث خطأ أثناء إضافة المادة.')
    finally:
        conn.close()
    return redirect(url_for('admin_subjects'))

@app.route('/admin/monitor')
def admin_monitor():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query = '''
        SELECT h.*, s.subject_name, d.dept_name, t.full_name as teacher_name
        FROM handouts h
        JOIN subjects s ON h.subject_id = s.id
        JOIN departments d ON h.dept_id = d.id
        JOIN teachers t ON h.teacher_id = t.id
        ORDER BY h.id DESC
    '''
    cursor.execute(query)
    logs = cursor.fetchall()
    conn.close()
    
    ai_message = None
    if len(logs) == 0:
        ai_message = "🔔 تنبيه ذكي: لا توجد أي مذكرات مرفوعة حتى الآن."
    elif len(logs) > 20:
        ai_message = "🔔 تنبيه ذكي: تم رفع عدد كبير من المذكرات مؤخرًا، تحقق من صحتها."

    return render_template('admin/monitor.html', logs=logs, ai_alert=ai_message)

@app.route('/admin/delete_subject/<int:id>')
def delete_subject(id):
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    try:
        cursor.execute(f"DELETE FROM course_structure WHERE subject_id = {query_placeholder}", (id,))
        cursor.execute(f"DELETE FROM subjects WHERE id = {query_placeholder}", (id,))
        conn.commit()
        flash('تم حذف المادة بنجاح.')
        ai_message = "🔔 تنبيه ذكي: تم حذف مادة من النظام، تأكد من تحديث الهيكل الدراسي."
    except Exception as e:
        flash('حدث خطأ أثناء الحذف.')
        ai_message = "🔔 تنبيه ذكي: حدث خطأ أثناء محاولة حذف المادة."
    finally:
        conn.close()
    
    return redirect(url_for('admin_subjects', ai_alert=ai_message))

@app.route('/admin/reports')
def admin_reports():
    if 'user_id' not in session or session['role'] != 'admin':
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    try:
        cursor.execute("SELECT COUNT(*) FROM teachers WHERE email != 'admin@musaid.edu.ly'")
        t_row = cursor.fetchone()
        t_count = t_row[0] if isinstance(t_row, tuple) else t_row["COUNT(*)"]
        
        cursor.execute('SELECT COUNT(*) FROM subjects')
        s_row = cursor.fetchone()
        s_count = s_row[0] if isinstance(s_row, tuple) else s_row["COUNT(*)"]
        
        cursor.execute('SELECT COUNT(*) FROM handouts')
        l_row = cursor.fetchone()
        l_count = l_row[0] if isinstance(l_row, tuple) else l_row["COUNT(*)"]
        
        cursor.execute('SELECT COUNT(*) FROM handouts WHERE dept_id = 1')
        it_row = cursor.fetchone()
        it_count = it_row[0] if isinstance(it_row, tuple) else it_row["COUNT(*)"]
        
        cursor.execute('SELECT COUNT(*) FROM handouts WHERE dept_id = 2')
        acc_row = cursor.fetchone()
        acc_count = acc_row[0] if isinstance(acc_row, tuple) else acc_row["COUNT(*)"]
    except Exception as e:
        t_count = s_count = l_count = it_count = acc_count = 0
    finally:
        conn.close()
    
    ai_message = None
    if l_count == 0:
        ai_message = "🔔 تنبيه ذكي: لم يتم رفع أي مذكرة هذا الأسبوع."
    elif acc_count < 2:
        ai_message = "🔔 تنبيه ذكي: نشاط قسم الإدارة منخفض، يرجى المتابعة."
    
    return render_template('admin/reports.html', 
                           t_count=t_count, 
                           s_count=s_count, 
                           l_count=l_count,
                           it_count=it_count,
                           acc_count=acc_count,
                           ai_alert=ai_message)

@app.route('/change_password', methods=['GET', 'POST'])
def change_password():
    if 'user_id' not in session:
        return redirect(url_for('login'))

    if request.method == 'POST':
        current = request.form.get('current_password') or ''
        new = request.form.get('new_password') or ''
        confirm = request.form.get('confirm_password') or ''

        conn = get_db_connection()
        cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
        query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
        
        cursor.execute(f'SELECT * FROM teachers WHERE id = {query_placeholder}', (session['user_id'],))
        user = cursor.fetchone()

        if not user or not verify_password(user['password'], current):
            conn.close()
            audit_log('password_change_failed', 'reason=wrong_current')
            flash('❌ كلمة المرور الحالية غير صحيحة.')
            return redirect(url_for('change_password'))

        if new != confirm:
            conn.close()
            flash('⚠️ كلمة المرور الجديدة وتأكيدها غير متطابقين.')
            return redirect(url_for('change_password'))

        ok, msg = validate_password(new)
        if not ok:
            conn.close()
            flash('⚠️ ' + msg)
            return redirect(url_for('change_password'))

        cursor.execute(f'UPDATE teachers SET password = {query_placeholder} WHERE id = {query_placeholder}',
                     (hash_password(new), session['user_id']))
        conn.commit()
        conn.close()
        audit_log('password_change_success')
        flash('✅ تم تغيير كلمة المرور بنجاح.')
        dest = 'admin_dashboard' if session.get('role') == 'admin' else 'teacher_dashboard'
        return redirect(url_for(dest))

    return render_template('change_password.html')

@app.route('/logout')
def logout():
    audit_log('logout')
    session.clear()
    flash('تم تسجيل الخروج بنجاح.')
    return redirect(url_for('login'))

# --- الدوال المساعدة للجافاسكريبت ---
@app.route('/get_semesters/<int:dept_id>')
def get_semesters(dept_id):
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    cursor.execute(f'''
        SELECT DISTINCT semester FROM course_structure 
        WHERE dept_id = {query_placeholder} ORDER BY semester
    ''', (dept_id,))
    semesters = cursor.fetchall()
    conn.close()
    if not semesters:
        return {"semesters": [{"id": i, "number": i} for i in range(1, 7)]}
    return {"semesters": [{"id": s['semester'], "number": s['semester']} for s in semesters]}

@app.route('/get_subjects/<int:dept_id>/<int:semester>')
def get_subjects(dept_id, semester):
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    try:
        cursor.execute(f'''
            SELECT s.id, s.subject_name FROM subjects s
            JOIN course_structure cs ON s.id = cs.subject_id
            WHERE cs.dept_id = {query_placeholder} AND cs.semester = {query_placeholder}
        ''', (dept_id, semester))
        subjects = cursor.fetchall()
    except:
        subjects = []
    conn.close()
    return {"subjects": [{"id": s['id'], "name": s['subject_name']} for s in subjects]}

@app.route('/delete_handout/<int:handout_id>')
def delete_handout(handout_id):
    if 'user_id' not in session or session['role'] not in ['teacher', 'admin']:
        return redirect(url_for('login'))
    
    conn = get_db_connection()
    cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
    query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
    
    try:
        cursor.execute(f'SELECT * FROM handouts WHERE id = {query_placeholder}', (handout_id,))
        handout = cursor.fetchone()
        if handout:
            if session['role'] == 'admin' or handout['teacher_id'] == session['user_id']:
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], handout['file_path'])
                if os.path.exists(file_path):
                    os.remove(file_path)
                cursor.execute(f'DELETE FROM handouts WHERE id = {query_placeholder}', (handout_id,))
                conn.commit()
                flash('تم حذف المذكرة بنجاح')
    except Exception as e:
        flash(f'خطأ: {str(e)}')
    finally:
        conn.close()
    return redirect(request.referrer or url_for('index'))

with app.app_context():
    try:
        conn = get_db_connection()
        cursor = conn.cursor() if hasattr(conn, 'cursor') else conn
        query_placeholder = '%s' if (os.environ.get('MUSAID_ENV') == 'production' or not hasattr(conn, 'execute')) else '?'
        
        cursor.execute(f"SELECT id FROM teachers WHERE email = 'admin@musaid.edu.ly'")
        existing = cursor.fetchone()
        if not existing:
            cursor.execute(
                f"INSERT INTO teachers (full_name, email, password) VALUES ({query_placeholder}, {query_placeholder}, {query_placeholder})",
                ('مدير النظام', 'admin@musaid.edu.ly', hash_password('33557799')))
            conn.commit()
        conn.close()
    except Exception as _e:
        print(f'⚠️  تخطّى تهيئة حساب المدير: {_e}')

# ==== تشغيل التطبيق ====
if __name__ == '__main__':
    port = int(os.environ.get('PORT', 5000))
    host = '0.0.0.0'  
    debug_mode = False
    
    print(f"Starting Musaid Portal on {host}:{port}")
    print(f"Environment: {ENV}")
    print(f"Database: {DATABASE}")
    print(f"Debug: {debug_mode}")
    
    app.run(host=host, port=port, debug=debug_mode)